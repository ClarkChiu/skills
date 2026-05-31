#!/usr/bin/env python3
"""Extract content + readable style from a user's .pptx, for path D (ingest).

Emits JSON to stdout: per-slide title/body/bullets/notes/image-count, plus the
theme colors and fonts python-pptx can read. We use this to rebuild the deck as HTML
in a matching style — we never redistribute the user's file.

Usage:
    python3 extract_pptx.py theirs.pptx > extracted.json

Needs python-pptx: `pip install python-pptx`. Absent → loud error (exit 2).
"""
import sys
import os
import json


def main() -> int:
    args = [a for a in sys.argv[1:] if not a.startswith("-")]
    if not args:
        print("usage: extract_pptx.py theirs.pptx", file=sys.stderr)
        return 2
    path = args[0]
    if not os.path.exists(path):
        print(f"no such file: {path}", file=sys.stderr)
        return 2

    try:
        from pptx import Presentation
        from pptx.util import Emu  # noqa: F401  (imported to confirm the package is real)
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx\n"
              "Or have the user paste their outline instead.", file=sys.stderr)
        return 2

    prs = Presentation(path)
    slides = []
    fonts = set()
    for idx, slide in enumerate(prs.slides, 1):
        title, bullets, images = None, [], 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                images += 1
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                for r in para.runs:
                    if r.font.name:
                        fonts.add(r.font.name)
                if not text:
                    continue
                if title is None and shape == slide.shapes.title:
                    title = text
                else:
                    bullets.append({"text": text, "level": para.level})
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({
            "n": idx, "title": title, "bullets": bullets,
            "images": images, "notes": notes,
        })

    out = {
        "slide_count": len(slides),
        "fonts_seen": sorted(fonts),
        "slides": slides,
        "_note": "Content + readable style only. Rebuild as HTML in a matching preset; "
                 "do not redistribute the source .pptx.",
    }
    print(json.dumps(out, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
